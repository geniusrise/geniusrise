from abc import ABC, abstractmethod
import csv
import json
from xml.etree import ElementTree
from PyPDF2 import PdfFileReader
from openpyxl import load_workbook
from docx import Document
from pptx import Presentation
from os.path import getmtime
from time import ctime
from pyth.plugins.rtf15.reader import Rtf15Reader
from pyth.plugins.plaintext.writer import PlaintextWriter
from bs4 import BeautifulSoup
import mimetypes
import os
from typing import Dict, Any


class TextExtractor(ABC):
    """Abstract base class for text extractors."""

    def __init__(self, extension: str):
        self.extension = extension

    @abstractmethod
    def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Process a file and extract text and metadata.

        This method must be overridden by subclasses.

        :param file_path: The path to the file to process.
        :return: A dictionary containing the extracted text and metadata.
        """
        raise NotImplementedError()


class PDFExtractor(TextExtractor):
    """Text extractor for PDF files."""

    def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Extract text and metadata from a PDF file."""
        with open(file_path, "rb") as file:
            pdf = PdfFileReader(file)
            text = ""
            for page in range(pdf.getNumPages()):
                text += pdf.getPage(page).extractText()
            return {"text": text, "metadata": pdf.getDocumentInfo(), "created_at": ctime(getmtime(file_path))}


class CSVExtractor(TextExtractor):
    """Text extractor for CSV files."""

    def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Extract text from a CSV file."""
        with open(file_path, "r") as file:
            reader = csv.reader(file)
            text = " ".join([" ".join(row) for row in reader])
            return {"text": text, "created_at": ctime(getmtime(file_path))}


class ExcelExtractor(TextExtractor):
    """Text extractor for Excel files."""

    def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Extract text from an Excel file."""
        wb = load_workbook(filename=file_path, read_only=True)
        text = ""
        for sheet in wb:
            for row in sheet:
                for cell in row:
                    text += str(cell.value) + " "
        return {"text": text, "created_at": ctime(getmtime(file_path))}


class DocExtractor(TextExtractor):
    """Text extractor for Word documents."""

    def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Extract text from a Word document."""
        doc = Document(file_path)
        text = " ".join([paragraph.text for paragraph in doc.paragraphs])
        return {"text": text, "created_at": ctime(getmtime(file_path))}


class JSONExtractor(TextExtractor):
    """Text extractor for JSON files."""

    def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Extract text from a JSON file."""
        with open(file_path, "r") as file:
            data = json.load(file)
            text = json.dumps(data)
            return {"text": text, "created_at": ctime(getmtime(file_path))}


class XMLExtractor(TextExtractor):
    """Text extractor for XML files."""

    def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Extract text from an XML file."""
        tree = ElementTree.parse(file_path)
        root = tree.getroot()
        text = ElementTree.tostring(root, encoding="utf8").decode("utf8")
        return {"text": text, "created_at": ctime(getmtime(file_path))}


class PPTExtractor(TextExtractor):
    """Text extractor for PowerPoint presentations."""

    def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Extract text from a PowerPoint presentation."""
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + " "
        return {"text": text, "created_at": ctime(getmtime(file_path))}


class TXTExtractor(TextExtractor):
    """Text extractor for text files."""

    def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Extract text from a text file."""
        with open(file_path, "r") as file:
            text = file.read()
        return {"text": text, "created_at": ctime(getmtime(file_path))}


class RTFExtractor(TextExtractor):
    """Text extractor for RTF files."""

    def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Extract text from an RTF file."""
        doc = Rtf15Reader.read(open(file_path, "rb"))
        text = PlaintextWriter.write(doc).getvalue()
        return {"text": text, "created_at": ctime(getmtime(file_path))}


class HTMLExtractor(TextExtractor):
    """Text extractor for HTML files."""

    def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Extract text from an HTML file."""
        with open(file_path, "r") as file:
            soup = BeautifulSoup(file.read(), "html.parser")
            text = soup.get_text()
        return {"text": text, "created_at": ctime(getmtime(file_path))}


class MarkdownExtractor(TextExtractor):
    """Text extractor for Markdown files."""

    def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Extract text from a Markdown file."""
        with open(file_path, "r") as file:
            text = file.read()
        return {"text": text, "created_at": ctime(getmtime(file_path))}


def extract(file_path: str) -> Dict[str, Any]:
    """Extract text and metadata from a file.

    This function determines the file type based on its extension
    and uses the appropriate extractor to process the file.

    :param file_path: The path to the file to process.
    :return: A dictionary containing the original file path and the extracted data.
    """
    # Get the MIME type of the file
    mime_type, _ = mimetypes.guess_type(file_path)
    # Get the file extension
    extension = os.path.splitext(file_path)[1]

    # Map MIME types to extractor classes
    extractor_map = {
        ".pdf": PDFExtractor,
        ".csv": CSVExtractor,
        ".xlsx": ExcelExtractor,
        ".doc": DocExtractor,
        ".docx": DocExtractor,
        ".json": JSONExtractor,
        ".xml": XMLExtractor,
        ".ppt": PPTExtractor,
        ".rtf": RTFExtractor,
        ".html": HTMLExtractor,
        ".md": MarkdownExtractor,
    }

    # Default to TXTExtractor for text files and unknown MIME types
    extractor_class = TXTExtractor

    # If the file extension is in the map, use the corresponding extractor
    if extension in extractor_map:
        extractor_class = extractor_map[extension]  # type: ignore

    # Create an instance of the extractor class and process the file
    extractor = extractor_class(extension)
    extracted_data = extractor.process(file_path)

    # Return both the original and extracted data
    return extracted_data
