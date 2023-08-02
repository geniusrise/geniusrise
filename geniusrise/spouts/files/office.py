# 🧠 Geniusrise
# Copyright (C) 2023  geniusrise.ai
#
# This program is free software: you can redistribute it and/or modify
# it under the terms of the GNU Affero General Public License as published by
# the Free Software Foundation, either version 3 of the License, or
# (at your option) any later version.
#
# This program is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
# GNU Affero General Public License for more details.
#
# You should have received a copy of the GNU Affero General Public License
# along with this program.  If not, see <http://www.gnu.org/licenses/>.

from os.path import getmtime
from time import ctime
from typing import Any, Dict

from docx import Document
from odf import teletype
from odf import text as odf_text
from odf.opendocument import load
from pptx import Presentation

from geniusrise.spouts.files.base import TextExtractor


class DocExtractor(TextExtractor):
    """Text extractor for Word documents."""

    def extract(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Extract text from a Word document."""
        doc = Document(file_path)
        text = " ".join([paragraph.text for paragraph in doc.paragraphs])
        return {"text": text, "created_at": ctime(getmtime(file_path))}


class PPTExtractor(TextExtractor):
    """Text extractor for PowerPoint presentations."""

    def extract(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Extract text from a PowerPoint presentation."""
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + " "
        return {"text": text, "created_at": ctime(getmtime(file_path))}


class ODTExtractor(TextExtractor):
    def process(self, file_path: str, **kwargs) -> Dict[str, Any]:
        doc = load(file_path)
        all_paragraphs = doc.getElementsByType(odf_text.P)
        text = " ".join(teletype.extractText(p) for p in all_paragraphs)
        return {"text": text, "created_at": ctime(getmtime(file_path))}
